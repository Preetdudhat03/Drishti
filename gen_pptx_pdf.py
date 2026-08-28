# -*- coding: utf-8 -*-
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DOCS_DIR = os.path.join(BASE_DIR, "docs")
os.makedirs(DOCS_DIR, exist_ok=True)

# =========================================================================
# 1. GENERATE 13_Drishti_SIH_2026_Presentation.pptx
# =========================================================================
print("[1/2] Generating 13_Drishti_SIH_2026_Presentation.pptx...")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color definitions
C_NAVY_DARK = RGBColor(9, 14, 23)
C_NAVY_CARD = RGBColor(15, 23, 42)
C_TEAL = RGBColor(14, 165, 233)
C_EMERALD = RGBColor(16, 185, 129)
C_AMBER = RGBColor(245, 158, 11)
C_CRIMSON = RGBColor(239, 68, 68)
C_WHITE = RGBColor(255, 255, 255)
C_MUTED = RGBColor(148, 163, 184)

slides_data = [
    {
        "title": "DRISHTI (दृष्टि)",
        "subtitle": "Explainable AI-Assisted Diabetic Retinopathy Screening & Clinical Decision Support System",
        "tagline": "Clinical Intelligence + Human Care | SIH 2026 PS-26038 | MathWorks",
        "points": [
            "Frontline AI Retinal Screening for Primary Health Centres (PHCs)",
            "Automated Multi-Factor Image Quality Gating (Sharpness, Exposure, FOV)",
            "ResNet-18 5-Class Severity Staging + Grad-CAM Neural Attention",
            "Dual-Persona Telemedicine & Immutable PostgreSQL Clinical Audit Trail",
            "District-Scale Telemedicine Modeling for 100,000+ Annual Rural Screenings"
        ]
    },
    {
        "title": "The Diabetic Blindness Crisis in Rural India",
        "subtitle": "Overwhelming Patient Volume Meets Extreme Specialist Scarcity",
        "tagline": "77+ Million Diabetic Patients vs ~20,000 Ophthalmologists",
        "points": [
            "Silent Pathology: Early non-proliferative retinopathy exhibits zero symptoms until vision loss is irreversible.",
            "Rural Access Deficit: Specialist density in rural districts is < 1:100,000; patients travel 50-100 km for routine exams.",
            "Screening Throughput Collapse: Manual screening cannot keep pace with annual national examination mandates.",
            "Delayed Treatment Window: Routine screening backlogs cause patients to miss the window for laser photocoagulation."
        ]
    },
    {
        "title": "Why Traditional Medical AI Fails in the Field",
        "subtitle": "The Triad of Failure: Black-Box Opacity, Quality Blindness, Disconnected Workflows",
        "tagline": "A standalone Python script is not a clinical platform",
        "points": [
            "Black-Box Opacity: Traditional CNNs output an uninterpretable probability number without anatomical evidence.",
            "Quality Blindness: Standard models perform inference on blurred/dark photos, outputting dangerous false negatives.",
            "Telemedicine Disconnection: Lack of role-based portals, offline synchronization, and clinician validation workflows.",
            "Regulatory & Trust Gap: Clinicians cannot adopt unverified AI decisions without human-in-the-loop oversight."
        ]
    },
    {
        "title": "Drishti: End-to-End System Architecture",
        "subtitle": "Gated, Explainable, and Clinician-Validated Telemedicine Platform",
        "tagline": "AI Assists. Doctor Decides. Complete Safety First.",
        "points": [
            "Pillar 1 — Image Quality Gate: Multi-factor Laplacian focus, illumination, and FOV scoring blocks ungradable inputs.",
            "Pillar 2 — Explainable Deep Learning: Fine-tuned ResNet-18 paired with Grad-CAM neural attention overlays.",
            "Pillar 3 — Human-in-the-Loop: Dual-persona mobile client syncing with Supabase PostgreSQL and specialist review queue.",
            "Pillar 4 — District Scaling: M/M/c discrete-event queuing simulation in MATLAB/Simulink for 120,000 annual patients."
        ]
    },
    {
        "title": "Image Quality Gating & Safety Engine",
        "subtitle": "Preventing False Negatives Before Neural Inference Begins",
        "tagline": "Composite Quality Score = 0.45 * Focus + 0.35 * Illumination + 0.20 * FOV",
        "points": [
            "Focus & Blur: Discrete Laplacian variance evaluated exclusively on the segmented retinal foreground mask.",
            "Illumination & Exposure: Evaluates mean luminance, 5th/95th percentiles, and over/underexposure penalties.",
            "Field of View (FOV): Otsu contour thresholding verifies >= 65% retinal coverage without edge occlusion.",
            "Tri-State Gating: GOOD (>=0.70 -> AI), BORDERLINE (0.45-0.69 -> CLAHE -> AI), UNGRADABLE (<0.45 -> Block AI & Recapture)."
        ]
    },
    {
        "title": "Adaptive Retinal Preprocessing & Structure Analysis",
        "subtitle": "Enhancing Contrast and Extracting Anatomical Landmarks",
        "tagline": "Exploratory Computer Vision Tools Supporting Decision Review",
        "points": [
            "Auto-Cropping: Automatically identifies non-zero bounding box to strip black background borders.",
            "Green Channel CLAHE: Contrast-limited adaptive histogram equalization (ClipLimit=0.02, 8x8 tiles) highlights microaneurysms.",
            "Blood Vessel Enhancement: Morphological top-hat filtering with disc structuring element (r=8) isolates vascular tree.",
            "Optic Disc Centroid: Regional peak red intensity localization provides anatomical reference for clinical review."
        ]
    },
    {
        "title": "Deep Learning Backbone — ResNet-18",
        "subtitle": "Residual Feature Extraction with Sub-50ms CPU Edge Inference",
        "tagline": "Pretrained on ImageNet-1k, Fine-Tuned on 3,662 Real APTOS Images",
        "points": [
            "Residual Learning: Skip connections (F(x) + x) eliminate vanishing gradients during deep feature backprop.",
            "Layer Topology: 4 residual stages (64, 128, 256, 512 channels) capturing multi-scale microvascular lesions.",
            "Compact Edge Footprint: 11.2M parameters, 44.8 MB state dict, optimized for rural tablets and low-power servers.",
            "Multi-Class Linear Head: Global Average Pooling (512x1x1) -> Linear Head outputting 5 clinical severity logits."
        ]
    },
    {
        "title": "Real APTOS Training & Class Imbalance Strategy",
        "subtitle": "Inverse Frequency Weighting & Quadratic Weighted Kappa Checkpointing",
        "tagline": "Strict 70% Train (2563) / 15% Val (550) / 15% Test (549) Split",
        "points": [
            "Severe Class Skew: Level 0 represents 49.3% while Level 3 represents only 5.3% of the dataset.",
            "Inverse Class Weighting: Normalized loss weights (W0=0.41, W1=1.98, W2=0.73, W3=3.80, W4=2.48) penalize rare class misses.",
            "Optimization: AdamW optimizer (LR=1e-4, Weight Decay=1e-2) with ReduceLROnPlateau scheduler on validation QWK.",
            "Safe Augmentation: Random horizontal flips, rotation (+-15 deg), and color jitter simulate variable field conditions."
        ]
    },
    {
        "title": "Rigorous Benchmark Validation (Held-Out Test Set)",
        "subtitle": "Scientifically Honest Evaluation on 549 Unseen Retinal Images",
        "tagline": "QWK = 0.870 | Accuracy = 76.87% | Specificity = 96.62% | Sensitivity = 82.14%",
        "points": [
            "5-Class Multiclass Agreement: Quadratic Weighted Kappa of 0.870 reflects high ordinal clinical concordance.",
            "Binary Referable Specificity: 96.62% (SIH Target: >85% | MET) — correctly clears 314 of 325 non-referable patients.",
            "Binary Referable Sensitivity: 82.14% (SIH Target: >90% | Fully reported) — identifies 184 of 224 referable cases.",
            "Discriminative Power: ROC AUC of 0.980 demonstrates near-ideal separability between healthy and diseased retinas."
        ]
    },
    {
        "title": "Explainable AI — Grad-CAM Neural Attention",
        "subtitle": "Making Deep Learning Clinically Verifiable for Ophthalmologists",
        "tagline": "Gradients backpropagated from layer4[1].conv2 feature maps",
        "points": [
            "Gradient Backpropagation: Computes gradients of predicted class score with respect to 512 conv feature maps.",
            "Global Importance Weights: Global average pooling computes alpha_k weights filtered through ReLU.",
            "Multi-Layer Saliency Overlay: 224x224 JET heatmap alpha-blended over fundus photo at alpha=0.45.",
            "Clinical Trust: Clinicians visually verify that the AI triggered on macular exudates rather than camera artifacts."
        ]
    },
    {
        "title": "Human-in-the-Loop Telemedicine & Audit Trail",
        "subtitle": "AI Assists Frontline Workers While Physicians Retain Authority",
        "tagline": "Validate, Override, or Reject with Immutable Audit Logging",
        "points": [
            "Health Worker Mode: Rapid patient registration, image quality assessment, AI screening, and bilingual PDF export.",
            "Ophthalmologist Portal: Priority-sorted telemedicine review queue with dual fundus / Grad-CAM inspection.",
            "Clinician Review Actions: One-tap 'Validate AI Result', 'Clinician Override' with mandatory notes, or 'Mark Ungradable'.",
            "Immutable PostgreSQL Audit: Every decision, reviewer identity, timestamp, and override note logged to audit_events."
        ]
    },
    {
        "title": "Field-Ready Mobile & Cloud Infrastructure",
        "subtitle": "Flutter Material 3, Riverpod 2.5, Supabase PostgreSQL, Offline Sync",
        "tagline": "Engineered for intermittent rural connectivity",
        "points": [
            "Cross-Platform Frontend: Flutter client delivering responsive mobile, tablet, and web clinician interfaces.",
            "Cloud Backend: Supabase PostgreSQL database with Row Level Security (RLS) policies and JWT authentication.",
            "Cloud Storage: Secure private fundus image buckets with signed URL access control.",
            "Offline Sync Queue: Local encrypted SQLite storage with SHA-256 integrity hashes auto-syncs when signal returns."
        ]
    },
    {
        "title": "District-Scale Telemedicine Simulation (MATLAB/Simulink)",
        "subtitle": "Queuing Model for 120,000 Annual Diabetic Screenings",
        "tagline": "Baseline (104% Overload) vs Drishti AI Triage (29% Utilization)",
        "points": [
            "Rural District Workload: 120,000 annual screenings across 300 days yields arrival rate lambda = 50 patients/hour.",
            "Baseline Failure: 2 ophthalmologists (capacity mu = 48/hr) experience rho = 104.2% utilization -> infinite queue backlog.",
            "Drishti AI Triage: ~72% cleared as routine; specialist arrival rate drops to lambda = 14/hr -> utilization rho = 29.2%.",
            "Queue Delay Reduction: Clinician review waiting time plummets from >48 hours to under 4.2 minutes."
        ]
    },
    {
        "title": "Innovation, Differentiation & Competitive Edge",
        "subtitle": "Why Drishti Outperforms Academic Models and Standalone Apps",
        "tagline": "The Complete Safety, Explainability, and Telemedicine Ecosystem",
        "points": [
            "Integrated Quality Gate: Blocks ungradable images before inference (absent in academic CNN benchmarks).",
            "Explainable Decision Support: Full Grad-CAM visual verification built directly into the doctor's review queue.",
            "Dual-Persona Workflow: Seamless handoff between rural health workers and district hospital ophthalmologists.",
            "District Capacity Modeling: Validated in MATLAB/Simulink to guarantee healthcare system throughput."
        ]
    },
    {
        "title": "Roadmap & Healthcare Impact Vision",
        "subtitle": "Scaling from Hackathon Prototype to National Deployment",
        "tagline": "Bringing Specialist Eye Care to Every Rural Primary Health Centre",
        "points": [
            "Phase 1 (Current MVP): 5-Class ResNet-18, Quality Gate, Grad-CAM, Flutter/Supabase, Simulink Queuing.",
            "Phase 2 (Q4 2026): High-resolution patch cropping (512x512) to exceed 92% sensitivity; edge TFLite offline inference.",
            "Phase 3 (2027): Multi-centric pilot trial across 10 rural Primary Health Centres in partnership with state health missions.",
            "Core Commitment: 'Drishti bridges India's rural healthcare gap with clinical intelligence and compassionate human care.'"
        ]
    }
]

blank_slide_layout = prs.slide_layouts[6]

for idx, sdata in enumerate(slides_data):
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background card
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_NAVY_DARK
    bg.line.fill.background()
    
    # Top header bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_NAVY_CARD
    top_bar.line.color.rgb = C_TEAL
    top_bar.line.width = Pt(1.5)
    
    # Slide Number & Title in top bar
    tf_top = top_bar.text_frame
    tf_top.word_wrap = True
    p1 = tf_top.paragraphs[0]
    p1.text = f"SLIDE {idx+1}: {sdata['title']}"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    
    p2 = tf_top.add_paragraph()
    p2.text = sdata['subtitle']
    p2.font.size = Pt(13)
    p2.font.color.rgb = C_TEAL
    
    # Tagline badge
    tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.733), Inches(0.5))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(14, 165, 233)
    tag_box.line.fill.background()
    tf_tag = tag_box.text_frame
    ptag = tf_tag.paragraphs[0]
    ptag.text = sdata['tagline']
    ptag.font.size = Pt(12)
    ptag.font.bold = True
    ptag.font.color.rgb = C_WHITE
    ptag.alignment = PP_ALIGN.CENTER
    
    # Main Content Box
    content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6), Inches(11.733), Inches(4.3))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = C_NAVY_CARD
    content_box.line.color.rgb = RGBColor(30, 41, 59)
    content_box.line.width = Pt(1.0)
    
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    for p_idx, pt in enumerate(sdata['points']):
        p = tf_content.paragraphs[0] if p_idx == 0 else tf_content.add_paragraph()
        p.text = f"•  {pt}"
        p.font.size = Pt(14)
        p.font.color.rgb = C_WHITE
        p.space_after = Pt(12)

pptx_path = os.path.join(DOCS_DIR, "13_Drishti_SIH_2026_Presentation.pptx")
prs.save(pptx_path)
print(f"[OK] Wrote 13_Drishti_SIH_2026_Presentation.pptx ({os.path.getsize(pptx_path)} bytes)")


# =========================================================================
# 2. GENERATE 02_Drishti_Master_Technical_Report.pdf
# =========================================================================
print("[2/2] Generating 02_Drishti_Master_Technical_Report.pdf...")
pdf_path = os.path.join(DOCS_DIR, "02_Drishti_Master_Technical_Report.pdf")

doc = SimpleDocTemplate(
    pdf_path,
    pagesize=letter,
    rightMargin=40,
    leftMargin=40,
    topMargin=40,
    bottomMargin=40
)

styles = getSampleStyleSheet()

# Custom styles
style_title = ParagraphStyle(
    'DocTitle',
    parent=styles['Normal'],
    fontName='Helvetica-Bold',
    fontSize=22,
    leading=26,
    textColor=colors.HexColor('#0F172A'),
    alignment=1, # Center
    spaceAfter=6
)

style_sub = ParagraphStyle(
    'DocSub',
    parent=styles['Normal'],
    fontName='Helvetica',
    fontSize=11,
    leading=15,
    textColor=colors.HexColor('#0EA5E9'),
    alignment=1,
    spaceAfter=15
)

style_h1 = ParagraphStyle(
    'DocH1',
    parent=styles['Heading1'],
    fontName='Helvetica-Bold',
    fontSize=14,
    leading=18,
    textColor=colors.HexColor('#0F172A'),
    spaceBefore=14,
    spaceAfter=6
)

style_h2 = ParagraphStyle(
    'DocH2',
    parent=styles['Heading2'],
    fontName='Helvetica-Bold',
    fontSize=11,
    leading=15,
    textColor=colors.HexColor('#0D9488'),
    spaceBefore=10,
    spaceAfter=4
)

style_body = ParagraphStyle(
    'DocBody',
    parent=styles['Normal'],
    fontName='Helvetica',
    fontSize=9,
    leading=13,
    textColor=colors.HexColor('#334155'),
    spaceAfter=6
)

style_table_header = ParagraphStyle(
    'TableHeader',
    parent=styles['Normal'],
    fontName='Helvetica-Bold',
    fontSize=8.5,
    leading=11,
    textColor=colors.white,
    alignment=1
)

style_table_cell = ParagraphStyle(
    'TableCell',
    parent=styles['Normal'],
    fontName='Helvetica',
    fontSize=8,
    leading=11,
    textColor=colors.HexColor('#1E293B')
)

story = []

# Title Banner
story.append(Paragraph("DRISHTI (दृष्टि)", style_title))
story.append(Paragraph("Explainable AI-Assisted Diabetic Retinopathy Screening & Decision Support System<br/><b>Smart India Hackathon 2026 | PS-26038 | MathWorks</b>", style_sub))
story.append(HRFlowable(width="100%", thickness=1.5, color=colors.HexColor('#0EA5E9'), spaceAfter=12))

# Section 1: Executive Overview
story.append(Paragraph("1. Executive Summary & Problem Breakdown", style_h1))
story.append(Paragraph("Drishti addresses India's diabetic blindness crisis by deploying an explainable, multi-factor safety-gated AI screening platform directly to rural Primary Health Centres (PHCs). With over 77 million diabetic patients and fewer than 20,000 ophthalmologists nationwide, automated triage of routine cases is mathematically essential to eliminate clinical backlog and prevent irreversible vision loss.", style_body))

# Module Status Matrix Table
status_data = [
    [Paragraph("Module / Capability", style_table_header), Paragraph("Implementation Status", style_table_header), Paragraph("Verified Test Metric / Evidence", style_table_header)],
    [Paragraph("Image Quality Gating", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("Laplacian Variance + Illumination + FOV Score >= 0.70", style_table_cell)],
    [Paragraph("5-Class DR Classification", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("ResNet-18 | 76.87% Accuracy | QWK = 0.870", style_table_cell)],
    [Paragraph("Referable DR Triage", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("Specificity = 96.62% | Sensitivity = 82.14% | AUC = 0.980", style_table_cell)],
    [Paragraph("Explainable AI (Grad-CAM)", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("layer4[1].conv2 Backpropagation Attention Overlays", style_table_cell)],
    [Paragraph("Human-in-the-Loop Telemedicine", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("Review Queue + Validate/Override Actions + Audit Trail", style_table_cell)],
    [Paragraph("District Telemedicine Scaling", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("MATLAB/Simulink M/M/c Model (120,000 Annual Patients)", style_table_cell)],
    [Paragraph("Offline Rural Sync Queue", style_table_cell), Paragraph("IMPLEMENTED", style_table_cell), Paragraph("SHA-256 Hash Idempotent SQLite / Storage Cache", style_table_cell)]
]

t_status = Table(status_data, colWidths=[140, 110, 280])
t_status.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0F172A')),
    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#F8FAFC'), colors.white]),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
]))
story.append(t_status)
story.append(Spacer(1, 10))

# Section 2: Real Model Evaluation
story.append(Paragraph("2. Deep Learning Validation on Real APTOS 2019 (Held-Out Test Set)", style_h1))
story.append(Paragraph("The model was evaluated strictly on a held-out test split of 549 images with zero data leakage. Multi-class Quadratic Weighted Kappa reached <b>0.870</b>, reflecting substantial agreement with international clinical standards. Binary referable specificity achieved <b>96.62%</b> with an ROC AUC of <b>0.980</b>. Measured sensitivity was <b>82.14%</b>, with errors concentrated in adjacent Grade 1 vs Grade 2 distinctions safely mitigated by mandatory clinician review.", style_body))

# 5-Class Confusion Matrix Table
cm_data = [
    [Paragraph("Ground Truth", style_table_header), Paragraph("Pred L0", style_table_header), Paragraph("Pred L1", style_table_header), Paragraph("Pred L2", style_table_header), Paragraph("Pred L3", style_table_header), Paragraph("Pred L4", style_table_header), Paragraph("Class F1", style_table_header)],
    [Paragraph("Level 0 (No DR)", style_table_cell), Paragraph("256", style_table_cell), Paragraph("12", style_table_cell), Paragraph("2", style_table_cell), Paragraph("0", style_table_cell), Paragraph("0", style_table_cell), Paragraph("96.8%", style_table_cell)],
    [Paragraph("Level 1 (Mild NPDR)", style_table_cell), Paragraph("3", style_table_cell), Paragraph("43", style_table_cell), Paragraph("9", style_table_cell), Paragraph("0", style_table_cell), Paragraph("0", style_table_cell), Paragraph("57.3%", style_table_cell)],
    [Paragraph("Level 2 (Moderate NPDR)", style_table_cell), Paragraph("0", style_table_cell), Paragraph("36", style_table_cell), Paragraph("92", style_table_cell), Paragraph("13", style_table_cell), Paragraph("9", style_table_cell), Paragraph("66.2%", style_table_cell)],
    [Paragraph("Level 3 (Severe NPDR)", style_table_cell), Paragraph("0", style_table_cell), Paragraph("0", style_table_cell), Paragraph("12", style_table_cell), Paragraph("9", style_table_cell), Paragraph("8", style_table_cell), Paragraph("31.6%", style_table_cell)],
    [Paragraph("Level 4 (Proliferative DR)", style_table_cell), Paragraph("0", style_table_cell), Paragraph("4", style_table_cell), Paragraph("13", style_table_cell), Paragraph("6", style_table_cell), Paragraph("22", style_table_cell), Paragraph("52.4%", style_table_cell)],
]
t_cm = Table(cm_data, colWidths=[120, 60, 60, 60, 60, 60, 110])
t_cm.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0D9488')),
    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#F0FDFA'), colors.white]),
    ('TOPPADDING', (0, 0), (-1, -1), 3),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
]))
story.append(t_cm)
story.append(Spacer(1, 10))

# Section 3: District Scaling & Telemedicine
story.append(Paragraph("3. District Telemedicine Capacity Simulation (MATLAB/Simulink)", style_h1))
story.append(Paragraph("Simulated 120,000 annual diabetic screenings across 300 working days (arrival rate lambda = 50 patients/hr). In the baseline manual system with 2 ophthalmologists (capacity mu = 48/hr), doctor utilization hits <b>104.2%</b>, causing infinite queue backlog and >48-hour delays. Drishti's automated AI triage clears ~72% of healthy patients, reducing specialist arrival rate to lambda = 14/hr (utilization <b>29.2%</b>) and cutting review delay to under <b>4.2 minutes</b>.", style_body))

# Build Document
doc.build(story)
print(f"[OK] Wrote 02_Drishti_Master_Technical_Report.pdf ({os.path.getsize(pdf_path)} bytes)")
print("Part 5 builders completed successfully.")
